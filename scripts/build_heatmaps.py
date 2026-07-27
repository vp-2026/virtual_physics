#!/usr/bin/env python3
import argparse
import copy
import csv
import json
import math
import random
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, Iterable, List, Optional, Sequence, Tuple

import matplotlib.pyplot as plt
import numpy as np
import pygame as pg
from matplotlib import colors as mcolors
try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    Presentation = RGBColor = MSO_AUTO_SHAPE_TYPE = PP_ALIGN = Inches = Pt = None

from make_trial import (
    ToolPicker,
    _apply_hidden_tool_behavior,
    circle_fully_in_polygon,
    current_goal_polygon,
    pick_container_or_none,
    point_in_poly,
    run_headless_episode,
)
from pyGameWorld.helpers import word2Color
from pyGameWorld.viewer import drawWorld
from tool_config import build_default_tools


FAMILY_ORDER = [
    "BackUp",
    "Balance",
    "BalanceUnder",
    "Basic",
    "Falling",
    "FallingAlt",
    "Prevention",
    "Remove",
]
NUMERIC_JSON_RE = re.compile(r"^(\d+)\.json$")
SLIDE_W = 13.333
SLIDE_H = 7.5


@dataclass
class EnvSpec:
    family: str
    env_num: int
    source_path: Path


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--env-root",
        default="data/environments",
    )
    parser.add_argument(
        "--asset-root",
        default="outputs/heatmaps",
    )
    parser.add_argument(
        "--output-ppt",
        default="outputs/environment_heatmaps.pptx",
    )
    parser.add_argument(
        "--gravity-mode",
        choices=("downward", "upward"),
        default="downward",
    )
    parser.add_argument("--seed", type=int, default=42)
    parser.add_argument("--grid-n", type=int, default=10)
    parser.add_argument("--samples-per-cell", type=int, default=2)
    parser.add_argument("--anchor-samples-per-cell", type=int, default=6)
    parser.add_argument("--candidate-points-per-cell", type=int, default=16)
    parser.add_argument("--top-cells", type=int, default=32)
    parser.add_argument("--fallback-random-points", type=int, default=240)
    parser.add_argument("--max-sim-seconds", type=float, default=12.0)
    parser.add_argument("--exact-heatmap", action="store_true")
    parser.add_argument("--heatmap-only", action="store_true")
    parser.add_argument("--redraw-from-csv", action="store_true")
    parser.add_argument(
        "--families",
        nargs="+",
        choices=FAMILY_ORDER,
        help="Limit processing to the listed environment families.",
    )
    parser.add_argument(
        "--skip-presentation",
        action="store_true",
        help="Do not rebuild the PowerPoint at the end of the run.",
    )
    parser.add_argument(
        "--skip-manifest-save",
        action="store_true",
        help="Do not write manifest.json during this run.",
    )
    parser.add_argument("--lattice-stride-px", type=int, default=2)
    return parser.parse_args()


def resolve_runtime_settings(args: argparse.Namespace) -> dict:
    settings = {
        "grid_n": args.grid_n,
        "samples_per_cell": args.samples_per_cell,
        "anchor_samples_per_cell": args.anchor_samples_per_cell,
        "candidate_points_per_cell": args.candidate_points_per_cell,
        "top_cells": args.top_cells,
        "fallback_random_points": args.fallback_random_points,
        "max_sim_seconds": args.max_sim_seconds,
        "exact_heatmap": args.exact_heatmap,
        "lattice_stride_px": args.lattice_stride_px,
    }
    if args.gravity_mode == "upward":
        if args.grid_n == 10:
            settings["grid_n"] = 8
        if args.samples_per_cell == 2:
            settings["samples_per_cell"] = 1
        if args.anchor_samples_per_cell == 6:
            settings["anchor_samples_per_cell"] = 8
        if args.candidate_points_per_cell == 16:
            settings["candidate_points_per_cell"] = 12
        if args.top_cells == 32:
            settings["top_cells"] = 24
        if args.fallback_random_points == 240:
            settings["fallback_random_points"] = 120
        if args.max_sim_seconds == 12.0:
            settings["max_sim_seconds"] = 10.0
    return settings


def orange_obj1_tools(gravity_mode: str = "downward") -> Dict[str, dict]:
    tools = build_default_tools("orange", "lightblue", "pink")
    if gravity_mode == "upward":
        tools["obj1"]["inverse_gravity"] = True
    return tools


def infer_target_ball_name(objects: Dict[str, dict]) -> str:
    for name, obj in objects.items():
        if obj.get("type") == "Ball" and str(obj.get("color", "")).lower() == "red":
            return name
    for name, obj in objects.items():
        if obj.get("type") == "Ball":
            return name
    return "Ball"


def make_runtime_world(source_path: Path, gravity_mode: str = "downward") -> dict:
    with open(source_path, "r", encoding="utf-8") as handle:
        world = json.load(handle)

    tools = orange_obj1_tools(gravity_mode=gravity_mode)
    world["tools"] = {name: copy.deepcopy(tool["polys"]) for name, tool in tools.items()}
    world.setdefault("world", {})
    world["world"].setdefault("objects", {})

    gcond = copy.deepcopy(world["world"].get("gcond") or {})
    objects = world["world"]["objects"]
    goal_name = gcond.get("goal")
    if goal_name not in objects:
        inferred_goal, _ = pick_container_or_none(objects, preferred_name=goal_name)
        goal_name = inferred_goal or "Goal"
    obj_name = gcond.get("obj")
    if obj_name not in objects:
        obj_name = infer_target_ball_name(objects)

    gcond["type"] = "SpecificInGoal"
    gcond["goal"] = goal_name
    gcond["obj"] = obj_name
    gcond["duration"] = str(gcond.get("duration", "3"))
    world["world"]["gcond"] = gcond
    try:
        gravity_mag = abs(float(world["world"].get("gravity", 200)))
    except Exception:
        gravity_mag = 200.0
    world["world"]["gravity"] = gravity_mag
    return world


def save_json(path: Path, payload: dict) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with open(path, "w", encoding="utf-8") as handle:
        json.dump(payload, handle, indent=2)


def render_world_rgb(world_dict: dict) -> np.ndarray:
    tp = ToolPicker(copy.deepcopy(world_dict))
    surf = drawWorld(tp._pyworld)
    arr = pg.surfarray.array3d(surf)
    return np.transpose(arr, (1, 0, 2))


def save_world_png(path: Path, world_dict: dict) -> None:
    bg = render_world_rgb(world_dict)
    path.parent.mkdir(parents=True, exist_ok=True)
    plt.imsave(path, bg)


def greens_alpha_cmap(alpha_max: float = 0.8, alpha_gamma: float = 0.55):
    base = plt.get_cmap("Greens")
    rgba = base(np.linspace(0, 1, 256))
    alpha_ramp = np.linspace(0.0, 1.0, 256) ** alpha_gamma
    rgba[:, 3] = alpha_ramp * alpha_max
    cmap = mcolors.ListedColormap(rgba)
    cmap.set_bad((0, 0, 0, 0))
    return cmap


def ensure_min_visible_score(scores: np.ndarray, min_visible: float = 0.03) -> np.ndarray:
    display_scores = np.array(scores, copy=True)
    mask = np.isfinite(display_scores) & (display_scores > 0.0)
    display_scores[mask] = np.maximum(display_scores[mask], min_visible)
    return display_scores


def is_still(pyworld, linear_threshold: float = 1.0, angular_threshold: float = 0.1) -> bool:
    for obj in pyworld.getDynamicObjects():
        if obj._cpBody.is_sleeping:
            continue
        lin_sq = obj._cpBody.velocity.length ** 2
        ang_v = abs(obj._cpBody.angular_velocity)
        if lin_sq > (linear_threshold ** 2) or ang_v > angular_threshold:
            return False
    return True


def goal_currently_satisfied(
    pyworld,
    world_dict: dict,
    target_obj_name: str,
    goal_name: str,
    goal_pts_world,
) -> bool:
    if goal_pts_world is None or target_obj_name not in pyworld.objects:
        return False

    target = pyworld.objects[target_obj_name]
    shape_def = world_dict["world"]["objects"].get(target_obj_name, {})
    live_goal_pts_world = current_goal_polygon(pyworld, goal_name, goal_pts_world)
    pos = target.getPos()
    if hasattr(pos, "x") and hasattr(pos, "y"):
        cx, cy = float(pos.x), float(pos.y)
    elif hasattr(pos, "tolist"):
        cx, cy = [float(value) for value in pos.tolist()[:2]]
    else:
        cx, cy = float(pos[0]), float(pos[1])

    if hasattr(target, "radius") or ("radius" in shape_def):
        radius = float(getattr(target, "radius", shape_def.get("radius", 0.0)))
        return (live_goal_pts_world is not None) and circle_fully_in_polygon(cx, cy, radius, live_goal_pts_world)

    poly_world = None
    if hasattr(target, "getVertices"):
        try:
            poly_world = [
                (float(v.x), float(v.y)) if hasattr(v, "x") else (float(v[0]), float(v[1]))
                for v in target.getVertices()
            ]
        except Exception:
            poly_world = None
    if poly_world is None and "vertices" in shape_def:
        poly_world = [(float(x), float(y)) for x, y in shape_def["vertices"]]
    return (live_goal_pts_world is not None) and (poly_world is not None) and all(
        point_in_poly(x, y, live_goal_pts_world) for x, y in poly_world
    )


def simulate_success(
    world_dict: dict,
    tool_def: dict,
    drop_xy: Tuple[int, int],
    max_sim_seconds: float,
    fps: float = 60.0,
    stillness_duration: float = 0.5,
) -> bool:
    sx, sy = int(drop_xy[0]), int(drop_xy[1])
    tp = ToolPicker(copy.deepcopy(world_dict))
    if tp._pyworld.checkCircleCollision((sx, sy), tool_def.get("placement_radius", 36)):
        return False

    translated_polys = []
    for poly in tool_def["polys"]:
        translated_polys.append([(vx + sx, vy + sy) for vx, vy in poly])

    tp._pyworld.addPlacedCompound(
        "PLACED",
        translated_polys,
        word2Color(tool_def["color"]),
        density=tool_def.get("density", 1.0),
        friction=tool_def.get("friction", 0.5),
        elasticity=tool_def.get("elasticity", 0.5),
    )
    _apply_hidden_tool_behavior(tp, "obj1", tool_def)

    gcond = world_dict["world"].get("gcond") or {}
    target_obj_name = gcond.get("obj", "Ball")
    goal_name = gcond.get("goal", "Goal")
    _, goal_obj = pick_container_or_none(world_dict["world"]["objects"], preferred_name=goal_name)
    goal_pts_world = goal_obj["points"] if goal_obj else None
    required_goal_duration = float(gcond.get("duration", 0.0) or 0.0)

    dt = 1.0 / fps
    in_goal_for = 0.0
    elapsed = 0.0
    still_for = 0.0
    while elapsed < max_sim_seconds:
        if goal_currently_satisfied(tp._pyworld, world_dict, target_obj_name, goal_name, goal_pts_world):
            in_goal_for += dt
            if in_goal_for >= required_goal_duration:
                return True
        else:
            in_goal_for = 0.0

        tp._pyworld.step(dt)
        elapsed += dt

        if is_still(tp._pyworld):
            still_for += dt
            if still_for >= stillness_duration and in_goal_for <= 0.0:
                break
        else:
            still_for = 0.0
    return False


def valid_drop(world_dict: dict, tool_def: dict, drop_xy: Tuple[int, int]) -> bool:
    sim_w, sim_h = world_dict["world"]["dims"]
    sx, sy = drop_xy
    radius = float(tool_def.get("placement_radius", 36))
    if sx < radius or sx > sim_w - radius or sy < radius or sy > sim_h - radius:
        return False
    tp = ToolPicker(copy.deepcopy(world_dict))
    return not tp._pyworld.checkCircleCollision((sx, sy), radius)


def numeric_env_specs(env_root: Path) -> List[EnvSpec]:
    specs: List[EnvSpec] = []
    families = [family for family in FAMILY_ORDER if (env_root / family).is_dir()]
    for family in families:
        for path in sorted((env_root / family).iterdir(), key=lambda p: int(NUMERIC_JSON_RE.match(p.name).group(1)) if NUMERIC_JSON_RE.match(p.name) else 10**9):
            match = NUMERIC_JSON_RE.fullmatch(path.name)
            if not match:
                continue
            specs.append(EnvSpec(family=family, env_num=int(match.group(1)), source_path=path))
    return specs


def filter_env_specs(specs: Sequence[EnvSpec], families: Optional[Sequence[str]]) -> List[EnvSpec]:
    if not families:
        return list(specs)
    allowed = set(families)
    return [spec for spec in specs if spec.family in allowed]


def object_bbox(spec: dict) -> Optional[Tuple[float, float, float, float]]:
    if spec.get("type") == "Ball" and "position" in spec and "radius" in spec:
        cx, cy = spec["position"]
        r = float(spec["radius"])
        return float(cx - r), float(cx + r), float(cy - r), float(cy + r)
    if "vertices" in spec:
        pts = spec["vertices"]
    elif "points" in spec:
        pts = spec["points"]
    elif "polys" in spec:
        pts = [pt for poly in spec["polys"] for pt in poly]
    else:
        return None
    xs = [float(p[0]) for p in pts]
    ys = [float(p[1]) for p in pts]
    if not xs or not ys:
        return None
    return min(xs), max(xs), min(ys), max(ys)


def movable_anchor_points(world_dict: dict, tool_radius: float) -> List[Tuple[int, int]]:
    anchors: List[Tuple[int, int]] = []
    sim_w, sim_h = world_dict["world"]["dims"]
    margin = max(18.0, tool_radius * 1.1)
    vertical_offsets = (
        margin,
        margin * 1.8,
        margin * 2.6,
        margin * 3.4,
    )
    for name, spec in world_dict["world"]["objects"].items():
        if name.startswith("_"):
            continue
        try:
            density = float(spec.get("density", 0))
        except Exception:
            density = 0.0
        if density <= 0:
            continue
        bbox = object_bbox(spec)
        if bbox is None:
            continue
        x0, x1, y0, y1 = bbox
        cx = 0.5 * (x0 + x1)
        cy = 0.5 * (y0 + y1)
        candidates = [
            (cx, cy),
            (x0 - margin, cy),
            (x1 + margin, cy),
        ]
        for offset in vertical_offsets:
            candidates.extend(
                [
                    (cx, y1 + offset),
                    (cx, y0 - offset),
                    (x0, y1 + offset),
                    (x1, y1 + offset),
                    (x0, y0 - offset),
                    (x1, y0 - offset),
                    (cx - 0.25 * (x1 - x0), y1 + offset),
                    (cx + 0.25 * (x1 - x0), y1 + offset),
                    (cx - 0.25 * (x1 - x0), y0 - offset),
                    (cx + 0.25 * (x1 - x0), y0 - offset),
                ]
            )
        for px, py in candidates:
            px = int(round(max(tool_radius, min(sim_w - tool_radius, px))))
            py = int(round(max(tool_radius, min(sim_h - tool_radius, py))))
            anchors.append((px, py))
    deduped: List[Tuple[int, int]] = []
    seen = set()
    for point in anchors:
        if point in seen:
            continue
        deduped.append(point)
        seen.add(point)
    return deduped


def generate_cell_points(
    x0: float,
    x1: float,
    y0: float,
    y1: float,
    base_count: int,
    anchor_bonus_count: int,
    anchors: Sequence[Tuple[int, int]],
    rng: random.Random,
) -> List[Tuple[int, int]]:
    cx = int(round((x0 + x1) / 2.0))
    cy = int(round((y0 + y1) / 2.0))
    points = [
        (cx, cy),
        (int(round(x0 + 0.25 * (x1 - x0))), int(round(y0 + 0.25 * (y1 - y0)))),
        (int(round(x0 + 0.75 * (x1 - x0))), int(round(y0 + 0.25 * (y1 - y0)))),
        (int(round(x0 + 0.25 * (x1 - x0))), int(round(y0 + 0.75 * (y1 - y0)))),
        (int(round(x0 + 0.75 * (x1 - x0))), int(round(y0 + 0.75 * (y1 - y0)))),
    ]

    cell_w = x1 - x0
    cell_h = y1 - y0
    margin_x = 0.45 * cell_w
    margin_y = 0.45 * cell_h
    local_anchors = [
        point for point in anchors
        if (x0 - margin_x) <= point[0] <= (x1 + margin_x) and (y0 - margin_y) <= point[1] <= (y1 + margin_y)
    ]
    local_anchors = sorted(local_anchors, key=lambda point: (point[0] - cx) ** 2 + (point[1] - cy) ** 2)

    for ax, ay in local_anchors[:3]:
        clamped = (
            int(round(min(max(ax, x0), x1))),
            int(round(min(max(ay, y0), y1))),
        )
        points.append(clamped)
        for jx, jy in ((-0.18, 0.0), (0.18, 0.0), (0.0, -0.18), (0.0, 0.18)):
            points.append((
                int(round(min(max(clamped[0] + jx * cell_w, x0), x1))),
                int(round(min(max(clamped[1] + jy * cell_h, y0), y1))),
            ))

    target_count = base_count + (anchor_bonus_count if local_anchors else 0)
    while len(points) < target_count:
        points.append((int(round(rng.uniform(x0, x1))), int(round(rng.uniform(y0, y1)))))

    deduped: List[Tuple[int, int]] = []
    seen = set()
    for point in points:
        if point in seen:
            continue
        deduped.append(point)
        seen.add(point)
    return deduped


def generate_lattice_points(
    x0: float,
    x1: float,
    y0: float,
    y1: float,
    stride_px: int,
) -> List[Tuple[int, int]]:
    points: List[Tuple[int, int]] = []
    start_x = int(math.ceil(x0))
    start_y = int(math.ceil(y0))
    if stride_px > 1:
        start_x = start_x + ((stride_px - (start_x % stride_px)) % stride_px)
        start_y = start_y + ((stride_px - (start_y % stride_px)) % stride_px)
    end_x = int(math.floor(x1))
    end_y = int(math.floor(y1))
    for px in range(start_x, end_x + 1, stride_px):
        for py in range(start_y, end_y + 1, stride_px):
            points.append((px, py))
    return points


def select_evenly_spaced_points(
    points: Sequence[Tuple[int, int]],
    limit: int,
) -> List[Tuple[int, int]]:
    if limit <= 0 or len(points) <= limit:
        return list(points)
    if limit == 1:
        return [points[len(points) // 2]]
    chosen: List[Tuple[int, int]] = []
    last_index = len(points) - 1
    for idx in range(limit):
        pick = int(round(idx * last_index / float(limit - 1)))
        chosen.append(points[pick])
    deduped: List[Tuple[int, int]] = []
    seen = set()
    for point in chosen:
        if point in seen:
            continue
        deduped.append(point)
        seen.add(point)
    return deduped


def build_heatmap(
    world_dict: dict,
    tool_def: dict,
    grid_n: int,
    samples_per_cell: int,
    anchor_samples_per_cell: int,
    max_sim_seconds: float,
    rng: random.Random,
    lattice_stride_px: int = 2,
    lattice_points_per_cell: int = 12,
    priority_points: Optional[Sequence[Tuple[int, int]]] = None,
    exact_heatmap: bool = False,
) -> Tuple[np.ndarray, List[dict], List[dict]]:
    sim_w, sim_h = world_dict["world"]["dims"]
    x_edges = np.linspace(0.0, float(sim_w), grid_n + 1)
    y_edges = np.linspace(0.0, float(sim_h), grid_n + 1)
    scores = np.full((grid_n, grid_n), np.nan, dtype=np.float32)
    cells: List[dict] = []
    sample_records: List[dict] = []
    priority_points = list(priority_points or [])
    if exact_heatmap:
        totals = np.zeros((grid_n, grid_n), dtype=np.int32)
        successes = np.zeros((grid_n, grid_n), dtype=np.int32)
        cell_w = float(sim_w) / float(grid_n)
        cell_h = float(sim_h) / float(grid_n)
        for px in range(0, int(sim_w), lattice_stride_px):
            for py in range(0, int(sim_h), lattice_stride_px):
                ix = min(grid_n - 1, int(px / cell_w))
                iy = min(grid_n - 1, int(py / cell_h))
                point = (int(px), int(py))
                is_valid = valid_drop(world_dict, tool_def, point)
                success = False
                if is_valid:
                    totals[iy, ix] += 1
                    if simulate_success(world_dict, tool_def, point, max_sim_seconds=max_sim_seconds):
                        successes[iy, ix] += 1
                        success = True
                sample_records.append(
                    {
                        "x": int(px),
                        "y": int(py),
                        "cell_ix": ix,
                        "cell_iy": iy,
                        "valid_drop": 1 if is_valid else 0,
                        "success": 1 if success else 0,
                    }
                )
        for ix in range(grid_n):
            for iy in range(grid_n):
                total = int(totals[iy, ix])
                if total == 0:
                    continue
                score = float(successes[iy, ix]) / float(total)
                x0, x1 = x_edges[ix], x_edges[ix + 1]
                y0, y1 = y_edges[iy], y_edges[iy + 1]
                scores[iy, ix] = score
                center = (int(round((x0 + x1) / 2.0)), int(round((y0 + y1) / 2.0)))
                cells.append(
                    {
                        "ix": ix,
                        "iy": iy,
                        "score": score,
                        "bounds": [float(x0), float(x1), float(y0), float(y1)],
                        "center": center,
                        "valid_points": [],
                    }
                )
        return scores, cells, sample_records

    for ix in range(grid_n):
        for iy in range(grid_n):
            x0, x1 = x_edges[ix], x_edges[ix + 1]
            y0, y1 = y_edges[iy], y_edges[iy + 1]
            lattice_points = generate_lattice_points(x0, x1, y0, y1, stride_px=lattice_stride_px)
            candidates = select_evenly_spaced_points(lattice_points, lattice_points_per_cell)
            for point in priority_points:
                if x0 <= point[0] <= x1 and y0 <= point[1] <= y1:
                    candidates.append((int(point[0]), int(point[1])))
            total = 0
            successes = 0
            valid_points: List[Tuple[int, int]] = []
            for point in candidates:
                is_valid = valid_drop(world_dict, tool_def, point)
                success = False
                if not is_valid:
                    sample_records.append(
                        {
                            "x": int(point[0]),
                            "y": int(point[1]),
                            "cell_ix": ix,
                            "cell_iy": iy,
                            "valid_drop": 0,
                            "success": 0,
                        }
                    )
                    continue
                total += 1
                valid_points.append(point)
                if simulate_success(world_dict, tool_def, point, max_sim_seconds=max_sim_seconds):
                    successes += 1
                    success = True
                sample_records.append(
                    {
                        "x": int(point[0]),
                        "y": int(point[1]),
                        "cell_ix": ix,
                        "cell_iy": iy,
                        "valid_drop": 1,
                        "success": 1 if success else 0,
                    }
                )
            if total == 0:
                continue
            score = float(successes) / float(total)
            scores[iy, ix] = score
            center = (int(round((x0 + x1) / 2.0)), int(round((y0 + y1) / 2.0)))
            cells.append(
                {
                    "ix": ix,
                    "iy": iy,
                    "score": score,
                    "bounds": [float(x0), float(x1), float(y0), float(y1)],
                    "center": center,
                    "valid_points": valid_points[:],
                }
            )
    return scores, cells, sample_records


def save_heatmap_samples_csv(path: Path, sample_records: Sequence[dict]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    fieldnames = ["x", "y", "cell_ix", "cell_iy", "valid_drop", "success"]
    with open(path, "w", encoding="utf-8", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=fieldnames)
        writer.writeheader()
        writer.writerows(sample_records)


def summarize_sample_records(sample_records: Sequence[dict]) -> dict:
    valid_count = 0
    success_count = 0
    success_points: List[Tuple[int, int]] = []
    for row in sample_records:
        if int(row["valid_drop"]) != 1:
            continue
        valid_count += 1
        if int(row["success"]) == 1:
            success_count += 1
            success_points.append((int(row["x"]), int(row["y"])))
    success_pct_valid = (100.0 * success_count / valid_count) if valid_count else 0.0
    return {
        "valid_count": valid_count,
        "success_count": success_count,
        "success_pct_valid": success_pct_valid,
        "success_points": success_points,
    }


def summarize_sample_csv(sample_csv_path: Path) -> dict:
    with open(sample_csv_path, "r", encoding="utf-8", newline="") as handle:
        reader = csv.DictReader(handle)
        return summarize_sample_records(list(reader))


def pick_success_point(success_points: Sequence[Tuple[int, int]]) -> Optional[Tuple[int, int]]:
    if not success_points:
        return None
    xs = [point[0] for point in success_points]
    ys = [point[1] for point in success_points]
    cx = float(sum(xs)) / float(len(xs))
    cy = float(sum(ys)) / float(len(ys))
    return min(success_points, key=lambda point: ((point[0] - cx) ** 2 + (point[1] - cy) ** 2, point[1], point[0]))


def format_success_pct(percent: Optional[float]) -> str:
    if percent is None:
        return ""
    if percent <= 0:
        return "0%"
    if percent < 0.1:
        return f"{percent:.3f}%"
    if percent < 1.0:
        return f"{percent:.2f}%"
    return f"{percent:.1f}%"


def load_scores_from_sample_csv(sample_csv_path: Path, grid_n: int) -> np.ndarray:
    totals = np.zeros((grid_n, grid_n), dtype=np.int32)
    successes = np.zeros((grid_n, grid_n), dtype=np.int32)
    with open(sample_csv_path, "r", encoding="utf-8", newline="") as handle:
        reader = csv.DictReader(handle)
        for row in reader:
            if int(row["valid_drop"]) != 1:
                continue
            ix = int(row["cell_ix"])
            iy = int(row["cell_iy"])
            totals[iy, ix] += 1
            if int(row["success"]) == 1:
                successes[iy, ix] += 1
    scores = np.full((grid_n, grid_n), np.nan, dtype=np.float32)
    mask = totals > 0
    scores[mask] = successes[mask] / totals[mask]
    return scores


def save_heatmap_png(
    path: Path,
    world_dict: dict,
    scores: np.ndarray,
    title: str,
    chosen_point: Optional[Tuple[int, int]] = None,
) -> None:
    bg = render_world_rgb(world_dict)
    display_scores = ensure_min_visible_score(np.flipud(scores))
    sim_w, sim_h = world_dict["world"]["dims"]
    fig, ax = plt.subplots(figsize=(6, 6), constrained_layout=True)
    ax.imshow(bg, extent=[0, sim_w, sim_h, 0], origin="upper", aspect="auto", interpolation="nearest")
    ax.imshow(
        display_scores,
        extent=[0, sim_w, sim_h, 0],
        origin="upper",
        aspect="auto",
        interpolation="nearest",
        cmap=greens_alpha_cmap(),
        norm=mcolors.PowerNorm(gamma=0.45, vmin=0.0, vmax=1.0),
    )
    ax.set_xlim(0, sim_w)
    ax.set_ylim(sim_h, 0)
    ax.axis("off")
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, dpi=180, bbox_inches="tight", pad_inches=0.02)
    plt.close(fig)


def choose_success_point(
    world_dict: dict,
    tool_def: dict,
    cells: List[dict],
    candidate_points_per_cell: int,
    top_cells: int,
    fallback_random_points: int,
    max_sim_seconds: float,
    rng: random.Random,
) -> Tuple[Optional[Tuple[int, int]], List[Tuple[int, int]]]:
    tried: List[Tuple[int, int]] = []
    seen = set()

    def try_point(point: Tuple[int, int]) -> bool:
        if point in seen:
            return False
        seen.add(point)
        if not valid_drop(world_dict, tool_def, point):
            return False
        tried.append(point)
        return simulate_success(world_dict, tool_def, point, max_sim_seconds=max_sim_seconds)

    ranked = sorted(cells, key=lambda item: (-item["score"], item["iy"], item["ix"]))
    anchors = movable_anchor_points(world_dict, float(tool_def.get("placement_radius", 36)))

    for cell in ranked[:top_cells]:
        x0, x1, y0, y1 = cell["bounds"]
        candidates = generate_cell_points(
            x0,
            x1,
            y0,
            y1,
            base_count=candidate_points_per_cell,
            anchor_bonus_count=max(2, candidate_points_per_cell // 2),
            anchors=anchors,
            rng=rng,
        )
        for point in candidates:
            if try_point(point):
                return point, tried

    sim_w, sim_h = world_dict["world"]["dims"]
    radius = int(round(tool_def.get("placement_radius", 36)))
    fallback_points = anchors[:]
    while len(fallback_points) < fallback_random_points:
        fallback_points.append((
            rng.randint(radius, int(sim_w - radius)),
            rng.randint(radius, int(sim_h - radius)),
        ))
    for point in fallback_points:
        if try_point(point):
            return point, tried

    return None, tried


def render_success_frames(
    world_dict: dict,
    tool_def: dict,
    drop_xy: Tuple[int, int],
    asset_dir: Path,
    basename: str,
    tool_name: str = "obj1",
    all_tools: Optional[Dict[str, dict]] = None,
) -> Tuple[Path, Path, dict]:
    asset_dir.mkdir(parents=True, exist_ok=True)
    if all_tools is None:
        gravity_mode = "upward" if tool_def.get("inverse_gravity") else "downward"
        all_tools = orange_obj1_tools(gravity_mode=gravity_mode)
        all_tools[tool_name] = tool_def
    else:
        all_tools = copy.deepcopy(all_tools)
        all_tools[tool_name] = tool_def
    frame_dir = asset_dir / f"{basename}_frames"
    payload = run_headless_episode(
        world_dict,
        tool_name=tool_name,
        tools_dict=all_tools,
        drop_xy=drop_xy,
        no_tool=False,
        frame_dir=str(frame_dir),
        frame_count=32,
    )
    placement = payload["placements"][0]
    if not placement.get("success"):
        raise RuntimeError("Recorded rollout was not successful.")

    payload_path = asset_dir / f"{basename}.json"
    save_json(payload_path, payload)
    return frame_dir, payload


def env_asset_dir(asset_root: Path, spec: EnvSpec) -> Path:
    return asset_root / spec.family / f"env_{spec.env_num}"


def fit_box(img_path: Path, max_w: float, max_h: float) -> Tuple[float, float]:
    from PIL import Image

    with Image.open(img_path) as image:
        w, h = image.size
    scale = min(max_w / w, max_h / h)
    return w * scale, h * scale


def recover_existing_entry(spec: EnvSpec, family_dir: Path) -> Optional[dict]:
    screenshot_path = family_dir / "environment.png"
    heatmap_path = family_dir / "heatmap.png"
    heatmap_samples_path = family_dir / "heatmap_samples.csv"
    runtime_world_path = family_dir / "runtime_world.json"
    if not (screenshot_path.exists() and heatmap_path.exists() and runtime_world_path.exists()):
        return None

    rollout_jsons = sorted(path for path in family_dir.glob("*.json") if path.name != "runtime_world.json")
    frames_dir = None
    chosen_point = None
    time_to_success = None
    success = False

    if rollout_jsons:
        payload_path = rollout_jsons[0]
        stem = payload_path.stem
        candidate_frames = family_dir / f"{stem}_frames"
        if candidate_frames.exists():
            frames_dir = candidate_frames
        with open(payload_path, "r", encoding="utf-8") as handle:
            payload = json.load(handle)
        placements = payload.get("placements") or []
        if placements:
            placement = placements[0]
            coords = placement.get("placement_coords")
            if isinstance(coords, list) and len(coords) == 2:
                chosen_point = [int(coords[0]), int(coords[1])]
            time_to_success = placement.get("time_to_success")
            success = bool(placement.get("success"))

    return {
        "family": spec.family,
        "env_num": spec.env_num,
        "source_path": str(spec.source_path),
        "runtime_world_path": str(runtime_world_path),
        "screenshot_path": str(screenshot_path),
        "heatmap_path": str(heatmap_path),
        "heatmap_samples_path": str(heatmap_samples_path) if heatmap_samples_path.exists() else None,
        "frames_dir": str(frames_dir) if frames_dir else None,
        "chosen_point": chosen_point,
        "tried_points": [],
        "time_to_success": time_to_success,
        "success": success,
    }


def add_textbox(slide, left, top, width, height, text, font_size=20, bold=False, color=(30, 30, 30), align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return tx


def add_title_bar(slide, title: str, subtitle: Optional[str] = None) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(0.75))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(23, 53, 87)
    band.line.fill.background()
    add_textbox(slide, Inches(0.45), Inches(0.16), Inches(8.0), Inches(0.3), title, font_size=28, bold=True, color=(255, 255, 255))
    if subtitle:
        add_textbox(slide, Inches(8.65), Inches(0.2), Inches(4.1), Inches(0.24), subtitle, font_size=13, color=(218, 232, 252), align=PP_ALIGN.RIGHT)


def first_frame_path(frames_dir: Optional[str]) -> Optional[Path]:
    if not frames_dir:
        return None
    frames = sorted(Path(frames_dir).glob("frame_*.png"))
    return frames[0] if frames else None


GRID_LEFTS = [Inches(0.3), Inches(2.863), Inches(5.426), Inches(7.989), Inches(10.552)]
GRID_TOPS = [Inches(0.85), Inches(4.165)]
GRID_W = Inches(2.483)
GRID_H = Inches(2.915)
LABEL_H = Inches(0.22)


def add_plain_title(slide, title: str) -> None:
    add_textbox(slide, Inches(0.3), Inches(0.08), Inches(12.73), Inches(0.5), title, font_size=24, bold=False, color=(0, 0, 0))


def iter_grid_slots():
    index = 0
    for top in GRID_TOPS:
        for left in GRID_LEFTS:
            yield index, left, top
            index += 1


def add_grid_label(slide, left, top, text: str, right_text: Optional[str] = None) -> None:
    if not right_text:
        add_textbox(slide, left, top + GRID_H, GRID_W, LABEL_H, text, font_size=12, color=(0, 0, 0), align=PP_ALIGN.CENTER)
        return
    add_textbox(slide, left, top + GRID_H, Inches(0.42), LABEL_H, text, font_size=12, color=(0, 0, 0), align=PP_ALIGN.LEFT)
    add_textbox(slide, left + Inches(0.44), top + GRID_H, GRID_W - Inches(0.44), LABEL_H, right_text, font_size=10, color=(58, 95, 58), align=PP_ALIGN.RIGHT)


def add_placeholder(slide, left, top, text: str) -> None:
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, GRID_W, GRID_H)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(243, 245, 249)
    box.line.color.rgb = RGBColor(208, 214, 223)
    add_textbox(slide, left + Inches(0.15), top + Inches(1.18), GRID_W - Inches(0.3), Inches(0.5), text, font_size=14, color=(88, 94, 104), align=PP_ALIGN.CENTER)


def upsert_manifest_entry(manifest: dict, entry: dict) -> None:
    for index, current in enumerate(manifest["entries"]):
        if current["family"] == entry["family"] and int(current["env_num"]) == int(entry["env_num"]):
            manifest["entries"][index] = entry
            return
    manifest["entries"].append(entry)


def build_presentation(output_ppt: Path, manifest: dict) -> None:
    if Presentation is None:
        raise RuntimeError("Optional presentation export requires python-pptx.")
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    entries_by_family: Dict[str, List[dict]] = {}
    for entry in manifest["entries"]:
        entries_by_family.setdefault(entry["family"], []).append(entry)

    for family in FAMILY_ORDER:
        entries = sorted(entries_by_family.get(family, []), key=lambda entry: entry["env_num"])
        if not entries:
            continue

        screenshot_slide = prs.slides.add_slide(blank)
        add_plain_title(screenshot_slide, f"{family} - Environments")
        for slot_idx, left, top in iter_grid_slots():
            if slot_idx >= len(entries):
                break
            entry = entries[slot_idx]
            screenshot_slide.shapes.add_picture(str(entry["screenshot_path"]), left, top, GRID_W, GRID_H)
            add_grid_label(screenshot_slide, left, top, str(entry["env_num"]))

        heatmap_slide = prs.slides.add_slide(blank)
        add_plain_title(heatmap_slide, f"{family} - Success Heatmaps")
        for slot_idx, left, top in iter_grid_slots():
            if slot_idx >= len(entries):
                break
            entry = entries[slot_idx]
            heatmap_slide.shapes.add_picture(str(entry["heatmap_path"]), left, top, GRID_W, GRID_H)
            add_grid_label(
                heatmap_slide,
                left,
                top,
                str(entry["env_num"]),
                right_text=format_success_pct(entry.get("success_pct_valid")),
            )

        rollout_slide = prs.slides.add_slide(blank)
        add_plain_title(rollout_slide, f"{family} - Successful Rollouts")
        for slot_idx, left, top in iter_grid_slots():
            if slot_idx >= len(entries):
                break
            entry = entries[slot_idx]
            label = str(entry["env_num"])
            frame_path = first_frame_path(entry.get("frames_dir"))
            if frame_path:
                rollout_slide.shapes.add_picture(str(frame_path), left, top, GRID_W, GRID_H)
            else:
                add_placeholder(rollout_slide, left, top, "No success found")
            add_grid_label(rollout_slide, left, top, label)

    output_ppt.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_ppt))


def main() -> None:
    args = parse_args()
    random.seed(args.seed)
    np.random.seed(args.seed)

    env_root = Path(args.env_root)
    asset_root = Path(args.asset_root)
    output_ppt = Path(args.output_ppt)
    if args.gravity_mode == "upward":
        if str(asset_root) == "outputs/heatmaps":
            asset_root = Path("outputs/heatmaps_upward")
        if str(output_ppt) == "outputs/environment_heatmaps.pptx":
            output_ppt = Path("outputs/environment_heatmaps_upward.pptx")
    runtime_settings = resolve_runtime_settings(args)
    specs = filter_env_specs(numeric_env_specs(env_root), args.families)
    tool_def = orange_obj1_tools(gravity_mode=args.gravity_mode)["obj1"]
    manifest_path = asset_root / "manifest.json"

    manifest = {
        "seed": args.seed,
        "grid_n": runtime_settings["grid_n"],
        "gravity_mode": args.gravity_mode,
        "runtime_settings": runtime_settings,
        "entries": [],
    }
    if manifest_path.exists():
        with open(manifest_path, "r", encoding="utf-8") as handle:
            existing_manifest = json.load(handle)
        if isinstance(existing_manifest, dict) and isinstance(existing_manifest.get("entries"), list):
            manifest["entries"] = existing_manifest["entries"]
            manifest["gravity_mode"] = existing_manifest.get("gravity_mode", args.gravity_mode)

    if args.redraw_from_csv:
        for entry in manifest["entries"]:
            sample_csv_path = entry.get("heatmap_samples_path")
            runtime_world_path = entry.get("runtime_world_path")
            heatmap_path = entry.get("heatmap_path")
            if not sample_csv_path or not runtime_world_path or not heatmap_path:
                continue
            with open(runtime_world_path, "r", encoding="utf-8") as handle:
                runtime_world = json.load(handle)
            sample_csv = Path(sample_csv_path)
            scores = load_scores_from_sample_csv(sample_csv, runtime_settings["grid_n"])
            stats = summarize_sample_csv(sample_csv)
            save_heatmap_png(
                Path(heatmap_path),
                runtime_world,
                scores,
                f"{entry['family']} env {entry['env_num']} | orange obj1 success",
            )
            entry["valid_count"] = stats["valid_count"]
            entry["success_count"] = stats["success_count"]
            entry["success_pct_valid"] = stats["success_pct_valid"]
            exact_point = pick_success_point(stats["success_points"])
            if exact_point is not None and (not entry.get("frames_dir") or not entry.get("success")):
                family_dir = Path(entry["runtime_world_path"]).parent
                basename = f"{entry['family'].lower()}_{int(entry['env_num'])}_orange_ball"
                frames_dir, payload = render_success_frames(
                    runtime_world,
                    tool_def,
                    exact_point,
                    family_dir,
                    basename,
                )
                placement = payload["placements"][0]
                entry["frames_dir"] = str(frames_dir)
                entry["chosen_point"] = [int(exact_point[0]), int(exact_point[1])]
                entry["time_to_success"] = placement.get("time_to_success")
                entry["success"] = placement.get("success", False)
            print(f"Redrew heatmap for {entry['family']}/{entry['env_num']}")
        if not args.skip_presentation:
            build_presentation(output_ppt, manifest)
            print(f"Wrote {output_ppt}")
        if not args.skip_manifest_save:
            save_json(manifest_path, manifest)
        return

    existing_keys = {
        (entry["family"], int(entry["env_num"])): entry
        for entry in manifest["entries"]
    }

    for spec in specs:
        family_dir = env_asset_dir(asset_root, spec)
        family_dir.mkdir(parents=True, exist_ok=True)
        key = (spec.family, spec.env_num)

        if args.heatmap_only:
            existing_entry = existing_keys.get(key) or recover_existing_entry(spec, family_dir)
            runtime_world_path = family_dir / "runtime_world.json"
            if runtime_world_path.exists():
                with open(runtime_world_path, "r", encoding="utf-8") as handle:
                    runtime_world = json.load(handle)
            else:
                runtime_world = make_runtime_world(spec.source_path, gravity_mode=args.gravity_mode)
                save_json(runtime_world_path, runtime_world)

            screenshot_path = family_dir / "environment.png"
            if not screenshot_path.exists():
                save_world_png(screenshot_path, runtime_world)

            rng = random.Random(args.seed * 100000 + spec.env_num * 97 + sum(ord(ch) for ch in spec.family))
            scores, cells, sample_records = build_heatmap(
                runtime_world,
                tool_def,
                grid_n=runtime_settings["grid_n"],
                samples_per_cell=runtime_settings["samples_per_cell"],
                anchor_samples_per_cell=runtime_settings["anchor_samples_per_cell"],
                max_sim_seconds=runtime_settings["max_sim_seconds"],
                rng=rng,
                lattice_stride_px=runtime_settings["lattice_stride_px"],
                exact_heatmap=runtime_settings["exact_heatmap"],
            )

            heatmap_path = family_dir / "heatmap.png"
            heatmap_samples_path = family_dir / "heatmap_samples.csv"
            stats = summarize_sample_records(sample_records)
            save_heatmap_png(
                heatmap_path,
                runtime_world,
                scores,
                f"{spec.family} env {spec.env_num} | orange obj1 success",
                chosen_point=None,
            )
            save_heatmap_samples_csv(heatmap_samples_path, sample_records)

            entry = {
                "family": spec.family,
                "env_num": spec.env_num,
                "source_path": str(spec.source_path),
                "runtime_world_path": str(runtime_world_path),
                "screenshot_path": str(screenshot_path),
                "heatmap_path": str(heatmap_path),
                "heatmap_samples_path": str(heatmap_samples_path),
                "frames_dir": existing_entry.get("frames_dir") if existing_entry else None,
                "chosen_point": existing_entry.get("chosen_point") if existing_entry else None,
                "tried_points": existing_entry.get("tried_points", []) if existing_entry else [],
                "time_to_success": existing_entry.get("time_to_success") if existing_entry else None,
                "success": existing_entry.get("success", False) if existing_entry else False,
                "valid_count": stats["valid_count"],
                "success_count": stats["success_count"],
                "success_pct_valid": stats["success_pct_valid"],
            }
            upsert_manifest_entry(manifest, entry)
            existing_keys[key] = entry
            if not args.skip_manifest_save:
                save_json(manifest_path, manifest)
            print(f"Rebuilt heatmap for {spec.family}/{spec.env_num}")
            continue

        if key in existing_keys:
            print(f"Reusing manifest entry for {spec.family}/{spec.env_num}")
            continue

        recovered = recover_existing_entry(spec, family_dir)
        if recovered:
            manifest["entries"].append(recovered)
            existing_keys[key] = recovered
            if not args.skip_manifest_save:
                save_json(manifest_path, manifest)
            print(f"Recovered existing assets for {spec.family}/{spec.env_num}")
            continue

        runtime_world = make_runtime_world(spec.source_path, gravity_mode=args.gravity_mode)
        runtime_world_path = family_dir / "runtime_world.json"
        save_json(runtime_world_path, runtime_world)

        screenshot_path = family_dir / "environment.png"
        save_world_png(screenshot_path, runtime_world)

        rng = random.Random(args.seed * 100000 + spec.env_num * 97 + sum(ord(ch) for ch in spec.family))
        scores, cells, sample_records = build_heatmap(
            runtime_world,
            tool_def,
            grid_n=runtime_settings["grid_n"],
            samples_per_cell=runtime_settings["samples_per_cell"],
            anchor_samples_per_cell=runtime_settings["anchor_samples_per_cell"],
            max_sim_seconds=runtime_settings["max_sim_seconds"],
            rng=rng,
            lattice_stride_px=runtime_settings["lattice_stride_px"],
            exact_heatmap=runtime_settings["exact_heatmap"],
        )

        stats = summarize_sample_records(sample_records)
        chosen_point = pick_success_point(stats["success_points"])
        tried_points = [chosen_point] if chosen_point is not None else []

        if chosen_point is None:
            chosen_point, tried_points = choose_success_point(
                runtime_world,
                tool_def,
                cells,
                candidate_points_per_cell=runtime_settings["candidate_points_per_cell"],
                top_cells=runtime_settings["top_cells"],
                fallback_random_points=runtime_settings["fallback_random_points"],
                max_sim_seconds=runtime_settings["max_sim_seconds"],
                rng=rng,
            )
        if chosen_point is None:
            if args.gravity_mode == "upward":
                retry_candidate_points = max(runtime_settings["candidate_points_per_cell"] + 6, int(runtime_settings["candidate_points_per_cell"] * 1.5))
                retry_top_cells = max(runtime_settings["top_cells"] + 8, runtime_settings["top_cells"])
                retry_fallback_random = max(runtime_settings["fallback_random_points"] + 80, int(runtime_settings["fallback_random_points"] * 1.5))
            else:
                retry_candidate_points = max(args.candidate_points_per_cell * 2, args.candidate_points_per_cell + 16)
                retry_top_cells = max(args.top_cells * 2, len(cells))
                retry_fallback_random = max(args.fallback_random_points * 3, args.fallback_random_points + 480)
            chosen_point, tried_points = choose_success_point(
                runtime_world,
                tool_def,
                cells,
                candidate_points_per_cell=retry_candidate_points,
                top_cells=retry_top_cells,
                fallback_random_points=retry_fallback_random,
                max_sim_seconds=runtime_settings["max_sim_seconds"],
                rng=rng,
            )

        heatmap_path = family_dir / "heatmap.png"
        heatmap_samples_path = family_dir / "heatmap_samples.csv"
        save_heatmap_png(
            heatmap_path,
            runtime_world,
            scores,
            f"{spec.family} env {spec.env_num} | orange obj1 success",
            chosen_point=chosen_point,
        )
        save_heatmap_samples_csv(heatmap_samples_path, sample_records)

        frames_dir = None
        placement = {}
        if chosen_point is not None:
            basename = f"{spec.family.lower()}_{spec.env_num}_orange_ball"
            frames_dir, payload = render_success_frames(
                runtime_world,
                tool_def,
                chosen_point,
                family_dir,
                basename,
            )
            placement = payload["placements"][0]
        else:
            print(f"Warning: no successful rollout point found for {spec.family}/{spec.env_num}")

        entry = {
            "family": spec.family,
            "env_num": spec.env_num,
            "source_path": str(spec.source_path),
            "runtime_world_path": str(runtime_world_path),
            "screenshot_path": str(screenshot_path),
            "heatmap_path": str(heatmap_path),
            "heatmap_samples_path": str(heatmap_samples_path),
            "frames_dir": str(frames_dir) if frames_dir else None,
            "chosen_point": list(chosen_point) if chosen_point else None,
            "tried_points": [list(point) for point in tried_points],
            "time_to_success": placement.get("time_to_success"),
            "success": placement.get("success", False),
            "valid_count": stats["valid_count"],
            "success_count": stats["success_count"],
            "success_pct_valid": stats["success_pct_valid"],
        }
        upsert_manifest_entry(manifest, entry)
        existing_keys[key] = entry
        if not args.skip_manifest_save:
            save_json(manifest_path, manifest)
        print(f"Built assets for {spec.family}/{spec.env_num}")

    if not args.skip_presentation:
        build_presentation(output_ppt, manifest)
        print(f"Wrote {output_ppt}")
    if not args.skip_manifest_save:
        save_json(manifest_path, manifest)


if __name__ == "__main__":
    main()
